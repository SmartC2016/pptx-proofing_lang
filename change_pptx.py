# -*- coding: utf-8 -*-
"""
This program is exchanging / setting the proofing language for any text-containing object
in a given Powerpoint Presentation.

My usecase:

I often review / add to existing Powerpoint presentations coming from colleagues all over the world.
Our company language is british english. However, very often the people do not change the proofing language
in the text elements, so that every english text is underlined and assumed wrong, if you for example have
swedish as proofing language. Often our presentations are 50 and more pages with title elements, sub-title
and text. To change all the elements to english is very cumbersome and for sure annoying.

Then there was the idea, to 'fix' this issue using the awesome Python programming language! :-) Yeah!

I use the great python-pptx module from Steve Canny - a million thanks to Steve!

This program allows you to select your pptx, select a language that should be set to all text-containing objects and
save your pptx under a new name.

if you have any questions, please don't hesitate to come back to me.

-----------

Copyright (C) <2017>  <Christian Hetmann>

This program is free software: you can redistribute it and/or modify
it under the terms of the GNU General Public License as published by
the Free Software Foundation, either version 3 of the License, or
(at your option) any later version.

This program is distributed in the hope that it will be useful,
but WITHOUT ANY WARRANTY; without even the implied warranty of
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
GNU General Public License for more details.

You should have received a copy of the GNU General Public License
along with this program.  If not, see <http://www.gnu.org/licenses/>.

You can contact me: chhe1970@gmail.com

-----------
"""
__author__ = "Christian Hetmann"


from pptx import Presentation
from pptx.enum.lang import MSO_LANGUAGE_ID


LICENSE = """
Copyright (C) <2017>  <Christian Hetmann>

This program is free software: you can redistribute it and/or modify
it under the terms of the GNU General Public License as published by
the Free Software Foundation, either version 3 of the License, or
(at your option) any later version.

This program is distributed in the hope that it will be useful,
but WITHOUT ANY WARRANTY; without even the implied warranty of
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
GNU General Public License for more details.

You should have received a copy of the GNU General Public License
along with this program.  If not, see <http://www.gnu.org/licenses/>.

You can contact me: chhe1970@gmail.com
"""

# todo create a tkinter window with the following features
# todo -- open file dialog: open a pptx file
# todo -- select a language from a dropdown menu that you want to set in the complete presentation
# todo ---- (default ENGLISH_UK, because that is my favourite :-) )
# todo -- create a little window for "logging" / showing what the program has found in the given pptx
# todo -- have a start button for execution
# todo -- have a save button to save the presentation
# todo create list with all existing languages in order to populate the dropdown -- SOLVED
# todo add licence file -- SOLVED


# select ENGLISH_UK as the new language to be set - this should be changed in the future to pick any language
new_language = MSO_LANGUAGE_ID.ENGLISH_UK

all_existing_lng = [MSO_LANGUAGE_ID.AFRIKAANS, MSO_LANGUAGE_ID.ALBANIAN, MSO_LANGUAGE_ID.AMHARIC,
                    MSO_LANGUAGE_ID.ARABIC, MSO_LANGUAGE_ID.ARABIC_ALGERIA,
                    MSO_LANGUAGE_ID.ARABIC_BAHRAIN, MSO_LANGUAGE_ID.ARABIC_EGYPT, MSO_LANGUAGE_ID.ARABIC_IRAQ,
                    MSO_LANGUAGE_ID.ARABIC_JORDAN, MSO_LANGUAGE_ID.ARABIC_KUWAIT, MSO_LANGUAGE_ID.ARABIC_LEBANON,
                    MSO_LANGUAGE_ID.ARABIC_LIBYA, MSO_LANGUAGE_ID.ARABIC_MOROCCO, MSO_LANGUAGE_ID.ARABIC_OMAN,
                    MSO_LANGUAGE_ID.ARABIC_QATAR, MSO_LANGUAGE_ID.ARABIC_SYRIA, MSO_LANGUAGE_ID.ARABIC_TUNISIA,
                    MSO_LANGUAGE_ID.ARABIC_UAE, MSO_LANGUAGE_ID.ARABIC_YEMEN, MSO_LANGUAGE_ID.ARMENIAN,
                    MSO_LANGUAGE_ID.ASSAMESE, MSO_LANGUAGE_ID.AZERI_CYRILLIC, MSO_LANGUAGE_ID.AZERI_LATIN,
                    MSO_LANGUAGE_ID.BASQUE, MSO_LANGUAGE_ID.BELGIAN_DUTCH, MSO_LANGUAGE_ID.BELGIAN_FRENCH,
                    MSO_LANGUAGE_ID.BENGALI, MSO_LANGUAGE_ID.BOSNIAN,
                    MSO_LANGUAGE_ID.BOSNIAN_BOSNIA_HERZEGOVINA_CYRILLIC,
                    MSO_LANGUAGE_ID.BOSNIAN_BOSNIA_HERZEGOVINA_LATIN, MSO_LANGUAGE_ID.BRAZILIAN_PORTUGUESE,
                    MSO_LANGUAGE_ID.BULGARIAN, MSO_LANGUAGE_ID.BURMESE, MSO_LANGUAGE_ID.BYELORUSSIAN,
                    MSO_LANGUAGE_ID.CATALAN, MSO_LANGUAGE_ID.CHEROKEE, MSO_LANGUAGE_ID.CHINESE_HONG_KONG_SAR,
                    MSO_LANGUAGE_ID.CHINESE_MACAO_SAR, MSO_LANGUAGE_ID.CHINESE_SINGAPORE,
                    MSO_LANGUAGE_ID.CROATIAN,
                    MSO_LANGUAGE_ID.CZECH, MSO_LANGUAGE_ID.DANISH, MSO_LANGUAGE_ID.DIVEHI,
                    MSO_LANGUAGE_ID.DUTCH, MSO_LANGUAGE_ID.EDO, MSO_LANGUAGE_ID.ENGLISH_AUS,
                    MSO_LANGUAGE_ID.ENGLISH_BELIZE, MSO_LANGUAGE_ID.ENGLISH_CANADIAN,
                    MSO_LANGUAGE_ID.ENGLISH_CARIBBEAN,
                    MSO_LANGUAGE_ID.ENGLISH_INDONESIA, MSO_LANGUAGE_ID.ENGLISH_IRELAND,
                    MSO_LANGUAGE_ID.ENGLISH_JAMAICA,
                    MSO_LANGUAGE_ID.ENGLISH_NEW_ZEALAND, MSO_LANGUAGE_ID.ENGLISH_PHILIPPINES,
                    MSO_LANGUAGE_ID.ENGLISH_SOUTH_AFRICA, MSO_LANGUAGE_ID.ENGLISH_TRINIDAD_TOBAGO,
                    MSO_LANGUAGE_ID.ENGLISH_UK,
                    MSO_LANGUAGE_ID.ENGLISH_US, MSO_LANGUAGE_ID.ENGLISH_ZIMBABWE, MSO_LANGUAGE_ID.ESTONIAN,
                    MSO_LANGUAGE_ID.FAEROESE, MSO_LANGUAGE_ID.FARSI, MSO_LANGUAGE_ID.FILIPINO,
                    MSO_LANGUAGE_ID.FINNISH,
                    MSO_LANGUAGE_ID.FRANCH_CONGO_DRC, MSO_LANGUAGE_ID.FRENCH, MSO_LANGUAGE_ID.FRENCH_CAMEROON,
                    MSO_LANGUAGE_ID.FRENCH_CANADIAN, MSO_LANGUAGE_ID.FRENCH_COTED_IVOIRE,
                    MSO_LANGUAGE_ID.FRENCH_HAITI,
                    MSO_LANGUAGE_ID.FRENCH_LUXEMBOURG, MSO_LANGUAGE_ID.FRENCH_MALI, MSO_LANGUAGE_ID.FRENCH_MONACO,
                    MSO_LANGUAGE_ID.FRENCH_MOROCCO, MSO_LANGUAGE_ID.FRENCH_REUNION,
                    MSO_LANGUAGE_ID.FRENCH_SENEGAL,
                    MSO_LANGUAGE_ID.FRENCH_WEST_INDIES, MSO_LANGUAGE_ID.FRISIAN_NETHERLANDS,
                    MSO_LANGUAGE_ID.FULFULDE,
                    MSO_LANGUAGE_ID.GAELIC_IRELAND, MSO_LANGUAGE_ID.GAELIC_SCOTLAND, MSO_LANGUAGE_ID.GALICIAN,
                    MSO_LANGUAGE_ID.GEORGIAN, MSO_LANGUAGE_ID.GERMAN, MSO_LANGUAGE_ID.GERMAN_AUSTRIA,
                    MSO_LANGUAGE_ID.GERMAN_LIECHTENSTEIN, MSO_LANGUAGE_ID.GERMAN_LUXEMBOURG,
                    MSO_LANGUAGE_ID.GREEK,
                    MSO_LANGUAGE_ID.GUARANI, MSO_LANGUAGE_ID.GUJARATI, MSO_LANGUAGE_ID.HAUSA,
                    MSO_LANGUAGE_ID.HAWAIIAN,
                    MSO_LANGUAGE_ID.HEBREW, MSO_LANGUAGE_ID.HINDI, MSO_LANGUAGE_ID.HUNGARIAN,
                    MSO_LANGUAGE_ID.IBIBIO,
                    MSO_LANGUAGE_ID.ICELANDIC, MSO_LANGUAGE_ID.IGBO, MSO_LANGUAGE_ID.INDONESIAN,
                    MSO_LANGUAGE_ID.INUKTITUT,
                    MSO_LANGUAGE_ID.ITALIAN, MSO_LANGUAGE_ID.JAPANESE, MSO_LANGUAGE_ID.KANNADA,
                    MSO_LANGUAGE_ID.KANURI,
                    MSO_LANGUAGE_ID.KASHMIRI, MSO_LANGUAGE_ID.KASHMIRI_DEVANAGARI, MSO_LANGUAGE_ID.KAZAKH,
                    MSO_LANGUAGE_ID.KHMER, MSO_LANGUAGE_ID.KIRGHIZ, MSO_LANGUAGE_ID.KONKANI,
                    MSO_LANGUAGE_ID.KOREAN,
                    MSO_LANGUAGE_ID.KYRGYZ, MSO_LANGUAGE_ID.LAO, MSO_LANGUAGE_ID.LATIN, MSO_LANGUAGE_ID.LATVIAN,
                    MSO_LANGUAGE_ID.LITHUANIAN, MSO_LANGUAGE_ID.MACEDONINAN_FYROM,
                    MSO_LANGUAGE_ID.MALAY_BRUNEI_DARUSSALAM,
                    MSO_LANGUAGE_ID.MALAYALAM, MSO_LANGUAGE_ID.MALAYSIAN, MSO_LANGUAGE_ID.MALTESE,
                    MSO_LANGUAGE_ID.MANIPURI, MSO_LANGUAGE_ID.MAORI, MSO_LANGUAGE_ID.MARATHI,
                    MSO_LANGUAGE_ID.MEXICAN_SPANISH, MSO_LANGUAGE_ID.MONGOLIAN, MSO_LANGUAGE_ID.NEPALI,
                    MSO_LANGUAGE_ID.NO_PROOFING, MSO_LANGUAGE_ID.NORWEGIAN_BOKMOL,
                    MSO_LANGUAGE_ID.NORWEGIAN_NYNORSK,
                    MSO_LANGUAGE_ID.ORIYA, MSO_LANGUAGE_ID.OROMO, MSO_LANGUAGE_ID.PASHTO,
                    MSO_LANGUAGE_ID.POLISH,
                    MSO_LANGUAGE_ID.PORTUGUESE, MSO_LANGUAGE_ID.PUNJABI, MSO_LANGUAGE_ID.QUECHUA_BOLIVIA,
                    MSO_LANGUAGE_ID.QUECHUA_ECUADOR, MSO_LANGUAGE_ID.QUECHUA_PERU, MSO_LANGUAGE_ID.RHAETO_ROMANIC,
                    MSO_LANGUAGE_ID.ROMANIAN, MSO_LANGUAGE_ID.ROMANIAN_MOLDOVA, MSO_LANGUAGE_ID.RUSSIAN,
                    MSO_LANGUAGE_ID.RUSSIAN_MOLDOVA, MSO_LANGUAGE_ID.SAMI_LAPPISH, MSO_LANGUAGE_ID.SANSKRIT,
                    MSO_LANGUAGE_ID.SEPEDI, MSO_LANGUAGE_ID.SERBIAN_BOSNIA_HERZEGOVINA_CYRILLIC,
                    MSO_LANGUAGE_ID.SERBIAN_BOSNIA_HERZEGOVINA_LATIN, MSO_LANGUAGE_ID.SERBIAN_CYRILLIC,
                    MSO_LANGUAGE_ID.SERBIAN_LATIN, MSO_LANGUAGE_ID.SESOTHO, MSO_LANGUAGE_ID.SIMPLIFIED_CHINESE,
                    MSO_LANGUAGE_ID.SINDHI, MSO_LANGUAGE_ID.SINDHI_PAKISTAN, MSO_LANGUAGE_ID.SINHALESE,
                    MSO_LANGUAGE_ID.SLOVAK, MSO_LANGUAGE_ID.SLOVENIAN, MSO_LANGUAGE_ID.SOMALI,
                    MSO_LANGUAGE_ID.SORBIAN,
                    MSO_LANGUAGE_ID.SPANISH, MSO_LANGUAGE_ID.SPANISH_ARGENTINA, MSO_LANGUAGE_ID.SPANISH_BOLIVIA,
                    MSO_LANGUAGE_ID.SPANISH_CHILE, MSO_LANGUAGE_ID.SPANISH_COLOMBIA,
                    MSO_LANGUAGE_ID.SPANISH_COSTA_RICA,
                    MSO_LANGUAGE_ID.SPANISH_DOMINICAN_REPUBLIC, MSO_LANGUAGE_ID.SPANISH_ECUADOR,
                    MSO_LANGUAGE_ID.SPANISH_EL_SALVADOR, MSO_LANGUAGE_ID.SPANISH_GUATEMALA,
                    MSO_LANGUAGE_ID.SPANISH_HONDURAS,
                    MSO_LANGUAGE_ID.SPANISH_MODERN_SORT, MSO_LANGUAGE_ID.SPANISH_NICARAGUA,
                    MSO_LANGUAGE_ID.SPANISH_PANAMA,
                    MSO_LANGUAGE_ID.SPANISH_PARAGUAY, MSO_LANGUAGE_ID.SPANISH_PERU,
                    MSO_LANGUAGE_ID.SPANISH_PUERTO_RICO,
                    MSO_LANGUAGE_ID.SPANISH_URUGUAY, MSO_LANGUAGE_ID.SPANISH_VENEZUELA, MSO_LANGUAGE_ID.SUTU,
                    MSO_LANGUAGE_ID.SWAHILI, MSO_LANGUAGE_ID.SWEDISH, MSO_LANGUAGE_ID.SWEDISH_FINLAND,
                    MSO_LANGUAGE_ID.SWISS_FRENCH, MSO_LANGUAGE_ID.SWISS_GERMAN, MSO_LANGUAGE_ID.SWISS_ITALIAN,
                    MSO_LANGUAGE_ID.SYRIAC, MSO_LANGUAGE_ID.TAJIK, MSO_LANGUAGE_ID.TAMAZIGHT,
                    MSO_LANGUAGE_ID.TAMAZIGHT_LATIN, MSO_LANGUAGE_ID.TAMIL, MSO_LANGUAGE_ID.TATAR,
                    MSO_LANGUAGE_ID.TELUGU,
                    MSO_LANGUAGE_ID.THAI, MSO_LANGUAGE_ID.TIBETAN, MSO_LANGUAGE_ID.TIGRIGNA_ERITREA,
                    MSO_LANGUAGE_ID.TIGRIGNA_ETHIOPIC, MSO_LANGUAGE_ID.TRADITIONAL_CHINESE,
                    MSO_LANGUAGE_ID.TSONGA,
                    MSO_LANGUAGE_ID.TSWANA, MSO_LANGUAGE_ID.TURKISH, MSO_LANGUAGE_ID.TURKMEN,
                    MSO_LANGUAGE_ID.UKRAINIAN,
                    MSO_LANGUAGE_ID.URDU, MSO_LANGUAGE_ID.UZBEK_CYRILLIC, MSO_LANGUAGE_ID.UZBEK_LATIN,
                    MSO_LANGUAGE_ID.VENDA, MSO_LANGUAGE_ID.VIETNAMESE, MSO_LANGUAGE_ID.WELSH,
                    MSO_LANGUAGE_ID.XHOSA,
                    MSO_LANGUAGE_ID.YI, MSO_LANGUAGE_ID.YIDDISH, MSO_LANGUAGE_ID.YORUBA, MSO_LANGUAGE_ID.ZULU]

input_file = 'test_pptx.pptx'
output_file = input_file[:-5] + '_modified.pptx'

def change_language(presentation, new_language):
    # iterate through all slides
    for slide_no, slide in enumerate(prs.slides):
        print(f'Working on SLIDE NO# {slide_no+1}')
        # iterate through all shapes/objects on one slide
        for shape in slide.shapes:
            # check if the shape/object has text (pictures e.g. don't have text)
            if shape.has_text_frame:
                # check for each paragraph of text for the actual shape/object
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.language_id != new_language:
                            # display the current language and new language
                            print(f'Slide {slide_no+1}, {shape.name}, from {run.font.language_id} --> {new_language}')
                            # set the 'new_language'
                            run.font.language_id = new_language
                        else:
                            print(f'Slide {slide_no+1}, shape {shape.name} is OK')
            else:
                print(f'Slide {slide_no+1}: The object "{shape.name}" has no text.')
        if slide_no < len(prs.slides)-1:
            print('--------- next slide ---------')
        else:
            print('******* Finished *******')
    return

# Open the presentation
prs = Presentation(input_file)

# save pptx with new filename
prs.save(output_file)
